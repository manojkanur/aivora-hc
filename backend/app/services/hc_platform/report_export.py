"""Layout-aware exporters for a StudioOutputDocument.

The on-screen report is a {title, subtitle, sections:[{layout, data, footnote}]}
document. These exporters render each section by its LAYOUT into a clean,
professional PDF / DOCX / PPTX / HTML - real tables, KPI tiles, colored badges,
and embedded chart images - so the export matches the on-screen deck instead of
flattening everything to text.
"""
from __future__ import annotations

import io
import re
from typing import Any

from . import report_charts

# McKinsey-style light palette (hex).
INK = "#0f172a"
MUTED = "#64748b"
FAINT = "#94a3b8"
LINE = "#e2e8f0"
PANEL = "#f8fafc"
BLUE = "#0060FF"
GOOD = "#059669"
WARN = "#b45309"
BAD = "#dc2626"

_SENT_COLOR = {"good": GOOD, "warning": WARN, "bad": BAD, "neutral": MUTED,
               "critical": BAD, "high": BAD, "medium": WARN, "low": MUTED}


# ── shared helpers ────────────────────────────────────────────────────────────

def _clean(s: Any) -> str:
    t = str(s or "")
    return (t.replace("—", "-").replace("–", "-").replace("’", "'")
             .replace("‘", "'").replace("“", '"').replace("”", '"')
             .replace("…", "...").replace(" ", " "))


def _strip_md(s: Any) -> str:
    """Drop markdown bold/italic markers for plain-text renderers."""
    t = _clean(s)
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"(?<!\*)\*(?!\*)(.+?)\*", r"\1", t)
    return t


def _sections(content: dict[str, Any]) -> list[dict[str, Any]]:
    secs = content.get("sections")
    return [s for s in secs if isinstance(s, dict)] if isinstance(secs, list) else []


def is_studio_document(content: dict[str, Any]) -> bool:
    """True when the draft holds a StudioOutputDocument (has typed sections)."""
    return bool(_sections(content))


def _doc_meta(content: dict[str, Any]) -> tuple[str, str, str]:
    title = _clean(content.get("title") or "Advisory Report")
    subtitle = _clean(content.get("subtitle") or "")
    studio = _clean(content.get("studio_name") or str(content.get("studio_id") or "").split(":")[-1].replace("-", " ").title())
    return title, subtitle, studio


def _section_title(section: dict[str, Any], index: int) -> str:
    t = _clean(section.get("title") or section.get("id") or f"Section {index+1}")
    return re.sub(r"^\s*\d+[.)]\s*", "", t)  # strip any leading "N."


# ── table extraction per layout (used by all formats) ─────────────────────────

def _kpi_rows(data: dict[str, Any]) -> list[list[str]]:
    items = data.get("items") or data.get("kpis") or []
    rows = [["Metric", "Value", "Detail"]]
    for it in items:
        if not isinstance(it, dict):
            continue
        val = f"{_clean(it.get('value'))}{(' ' + _clean(it.get('unit'))) if it.get('unit') else ''}".strip()
        rows.append([_strip_md(it.get("label")), val, _strip_md(it.get("sublabel"))])
    return rows


def _comparison_rows(data: dict[str, Any]) -> list[list[str]]:
    cols = data.get("columns") or []
    keys = [c.get("key") for c in cols if isinstance(c, dict)]
    header = [_clean(c.get("label") or c.get("key")) for c in cols if isinstance(c, dict)]
    rows = [header] if header else []
    for r in (data.get("rows") or []):
        if isinstance(r, dict) and keys:
            rows.append([_strip_md(r.get(k, "")) for k in keys])
        elif isinstance(r, dict) and "cells" in r:
            cells = [_strip_md((c.get("value") if isinstance(c, dict) else c)) for c in (r.get("cells") or [])]
            rows.append([_strip_md(r.get("label", ""))] + cells)
    return rows


def _recommendation_rows(data: dict[str, Any]) -> list[list[str]]:
    rows = [["Recommendation", "Priority", "Rationale"]]
    for it in (data.get("items") or data.get("cards") or []):
        if not isinstance(it, dict):
            continue
        rows.append([_strip_md(it.get("title")), _clean(it.get("priority") or "").title(), _strip_md(it.get("rationale") or it.get("description"))])
    return rows


def _risk_rows(data: dict[str, Any]) -> list[list[str]]:
    rows = [["Risk", "Severity", "Mitigation"]]
    for it in (data.get("items") or []):
        if not isinstance(it, dict):
            continue
        rows.append([_strip_md(it.get("title")), _clean(it.get("severity") or "").title(), _strip_md(it.get("mitigation") or it.get("description"))])
    return rows


def _timeline_rows(data: dict[str, Any]) -> list[list[str]]:
    rows = [["Phase", "Action", "Owner"]]
    for h in (data.get("horizons") or []):
        if not isinstance(h, dict):
            continue
        label = _clean(h.get("label"))
        actions = h.get("actions") or []
        if not actions:
            rows.append([label, "", ""])
        for i, a in enumerate(actions):
            if isinstance(a, dict):
                rows.append([label if i == 0 else "", _strip_md(a.get("title")), _strip_md(a.get("owner"))])
            else:
                rows.append([label if i == 0 else "", _strip_md(a), ""])
    return rows


def _table_for(layout: str, data: dict[str, Any]) -> list[list[str]] | None:
    if layout in ("kpi_grid", "kpi-grid", "kpis", "stat_block", "stats"):
        return _kpi_rows(data)
    if layout in ("comparison_table", "table"):
        return _comparison_rows(data)
    if layout in ("recommendation_cards", "recommendations", "recommendation_grid"):
        return _recommendation_rows(data)
    if layout in ("risk_flags_list", "risk_flags", "risks"):
        return _risk_rows(data)
    if layout in ("timeline", "roadmap"):
        return _timeline_rows(data)
    return None


def _narrative_text(data: dict[str, Any]) -> str:
    body = data.get("body") or data.get("text") or ""
    if body:
        return _strip_md(body)
    # Hero layout: headline + subheadline + highlight bullets.
    headline = data.get("headline") or data.get("title") or ""
    sub = data.get("subheadline") or data.get("subtitle") or ""
    highlights = data.get("highlights")
    parts: list[str] = []
    if headline:
        parts.append(_strip_md(str(headline)))
    if sub:
        parts.append(_strip_md(str(sub)))
    if isinstance(highlights, list) and highlights:
        for h in highlights:
            if h:
                parts.append("- " + _strip_md(str(h)))
    return "\n".join(parts)


def _callout_text(data: dict[str, Any]) -> tuple[str, str]:
    return _strip_md(data.get("quote") or data.get("text") or data.get("body") or ""), _strip_md(data.get("attribution") or "")


def _footnote_url(section: dict[str, Any]) -> str:
    fn = _clean(section.get("footnote") or "")
    m = re.search(r"https?://\S+", fn)
    return m.group(0) if m else ""


# ── PDF (reportlab) ───────────────────────────────────────────────────────────

def build_pdf(content: dict[str, Any], org_name: str = "") -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
    )

    title, subtitle, studio = _doc_meta(content)
    ink = colors.HexColor(INK)
    blue = colors.HexColor(BLUE)
    line = colors.HexColor(LINE)
    panel = colors.HexColor(PANEL)
    muted = colors.HexColor(MUTED)

    ss = getSampleStyleSheet()
    h_title = ParagraphStyle("t", parent=ss["Title"], fontName="Helvetica-Bold", fontSize=24, leading=28, textColor=ink, spaceAfter=6)
    h_sub = ParagraphStyle("s", parent=ss["Normal"], fontSize=11, leading=15, textColor=muted, spaceAfter=2)
    h_eyebrow = ParagraphStyle("e", parent=ss["Normal"], fontName="Helvetica-Bold", fontSize=8.5, textColor=blue, spaceAfter=4)
    h_sec = ParagraphStyle("h2", parent=ss["Heading2"], fontName="Helvetica-Bold", fontSize=14, leading=17, textColor=ink, spaceBefore=14, spaceAfter=6)
    body = ParagraphStyle("b", parent=ss["Normal"], fontSize=10, leading=15, textColor=ink, alignment=TA_LEFT, spaceAfter=6)
    cap = ParagraphStyle("c", parent=ss["Normal"], fontSize=8, leading=11, textColor=muted, spaceAfter=8)
    cell = ParagraphStyle("cell", parent=ss["Normal"], fontSize=8.5, leading=12, textColor=ink)
    cell_h = ParagraphStyle("cellh", parent=ss["Normal"], fontName="Helvetica-Bold", fontSize=8.5, leading=12, textColor=colors.white)
    quote = ParagraphStyle("q", parent=ss["Normal"], fontName="Helvetica-Bold", fontSize=13, leading=18, textColor=ink, spaceAfter=4)

    story: list[Any] = []

    # Cover block
    story.append(Paragraph(studio.upper(), h_eyebrow))
    story.append(Paragraph(title, h_title))
    if subtitle:
        story.append(Paragraph(subtitle, h_sub))
    story.append(Spacer(1, 4))
    meta = f"Aivora HC Advisory   |   {len(_sections(content))} sections   |   Confidential"
    story.append(Paragraph(meta, cap))
    story.append(Table([[""]], colWidths=[170 * mm], style=TableStyle([("LINEBELOW", (0, 0), (-1, -1), 1.2, ink)])))
    story.append(Spacer(1, 8))

    def _styled_table(rows: list[list[str]]) -> Table:
        header, *data_rows = rows
        n = len(header)
        ncols = [Paragraph(_clean(h), cell_h) for h in header]
        prows = [ncols] + [[Paragraph(_clean(c), cell) for c in r] for r in data_rows]
        avail = 170 * mm
        # First column a bit wider for label tables.
        widths = [avail * (0.34 if n <= 3 else 0.28)] + [(avail * (0.66 if n <= 3 else 0.72)) / (n - 1)] * (n - 1) if n > 1 else [avail]
        t = Table(prows, colWidths=widths, repeatRows=1)
        style = [
            ("BACKGROUND", (0, 0), (-1, 0), blue),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("LINEBELOW", (0, 0), (-1, -1), 0.5, line),
            ("LINEBEFORE", (0, 0), (-1, -1), 0, line),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, panel]),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("LEFTPADDING", (0, 0), (-1, -1), 7),
            ("RIGHTPADDING", (0, 0), (-1, -1), 7),
            ("BOX", (0, 0), (-1, -1), 0.5, line),
        ]
        return Table(prows, colWidths=widths, repeatRows=1, style=TableStyle(style))

    for i, section in enumerate(_sections(content)):
        layout = str(section.get("layout") or "")
        data = section.get("data") if isinstance(section.get("data"), dict) else {}
        story.append(Paragraph(f"{i+1:02d}   {_section_title(section, i)}", h_sec))

        chart = report_charts.chart_png_for(layout, data)
        if chart:
            try:
                img = Image(io.BytesIO(chart))
                max_w = 165 * mm
                iw, ih = img.drawWidth, img.drawHeight
                if iw > max_w:
                    img.drawHeight = ih * (max_w / iw)
                    img.drawWidth = max_w
                story.append(img)
                story.append(Spacer(1, 4))
            except Exception:  # noqa: BLE001
                chart = None

        if layout in ("narrative_paragraph", "narrative", "prose", "text"):
            for para in _narrative_text(data).split("\n"):
                if para.strip():
                    story.append(Paragraph(_clean(para), body))
        elif layout in ("callout_quote", "quote", "callout", "note", "warning"):
            q, attr = _callout_text(data)
            story.append(Table([[Paragraph(q, quote)] + ([] if not attr else [])],
                               colWidths=[170 * mm],
                               style=TableStyle([("BACKGROUND", (0, 0), (-1, -1), panel),
                                                 ("LINEBEFORE", (0, 0), (0, -1), 3, blue),
                                                 ("LEFTPADDING", (0, 0), (-1, -1), 12),
                                                 ("TOPPADDING", (0, 0), (-1, -1), 10),
                                                 ("BOTTOMPADDING", (0, 0), (-1, -1), 10)])))
            if attr:
                story.append(Paragraph(attr, cap))
        else:
            rows = _table_for(layout, data)
            if rows and len(rows) > 1:
                story.append(_styled_table(rows))
            elif not chart:
                # Unknown/other layout: fall back to any narrative body.
                nb = _narrative_text(data)
                if nb:
                    story.append(Paragraph(nb, body))

        url = _footnote_url(section)
        if url:
            story.append(Paragraph(f'Source: <a href="{url}" color="#1d4ed8">{url}</a>', cap))
        story.append(Spacer(1, 4))

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=20 * mm, rightMargin=20 * mm,
                            topMargin=18 * mm, bottomMargin=18 * mm, title=title)

    def _footer(canvas, doc_):
        canvas.saveState()
        canvas.setFont("Helvetica", 7.5)
        canvas.setFillColor(muted)
        canvas.drawString(20 * mm, 10 * mm, title[:80])
        canvas.drawRightString(190 * mm, 10 * mm, f"Aivora HC   |   Page {doc_.page}")
        canvas.setStrokeColor(line)
        canvas.line(20 * mm, 13 * mm, 190 * mm, 13 * mm)
        canvas.restoreState()

    doc.build(story, onFirstPage=_footer, onLaterPages=_footer)
    return buf.getvalue()


# ── HTML (self-contained one-pager) ───────────────────────────────────────────

def build_html(content: dict[str, Any], org_name: str = "") -> bytes:
    import base64
    import html as _html

    title, subtitle, studio = _doc_meta(content)
    e = lambda s: _html.escape(_clean(s))

    def _table_html(rows: list[list[str]]) -> str:
        header, *body_rows = rows
        th = "".join(f"<th>{e(h)}</th>" for h in header)
        trs = "".join("<tr>" + "".join(f"<td>{e(c)}</td>" for c in r) + "</tr>" for r in body_rows)
        return f'<table><thead><tr>{th}</tr></thead><tbody>{trs}</tbody></table>'

    parts = []
    for i, section in enumerate(_sections(content)):
        layout = str(section.get("layout") or "")
        data = section.get("data") if isinstance(section.get("data"), dict) else {}
        parts.append(f'<section><h2><span class="num">{i+1:02d}</span> {e(_section_title(section, i))}</h2>')
        chart = report_charts.chart_png_for(layout, data)
        if chart:
            b64 = base64.b64encode(chart).decode()
            parts.append(f'<img class="chart" src="data:image/png;base64,{b64}" alt="chart"/>')
        if layout in ("narrative_paragraph", "narrative", "prose", "text", "hero"):
            for para in _narrative_text(data).split("\n"):
                p = para.strip()
                if not p:
                    continue
                if p.startswith("- "):
                    parts.append(f"<p style='margin-left:14px'>&bull; {e(p[2:])}</p>")
                else:
                    parts.append(f"<p>{e(p)}</p>")
        elif layout in ("callout_quote", "quote", "callout", "note", "warning"):
            q, attr = _callout_text(data)
            parts.append(f'<blockquote><p>{e(q)}</p>{f"<cite>{e(attr)}</cite>" if attr else ""}</blockquote>')
        else:
            rows = _table_for(layout, data)
            if rows and len(rows) > 1:
                parts.append(_table_html(rows))
            elif not chart:
                nb = _narrative_text(data)
                if nb:
                    parts.append(f"<p>{e(nb)}</p>")
        url = _footnote_url(section)
        if url:
            parts.append(f'<p class="src">Source: <a href="{e(url)}">{e(url)}</a></p>')
        parts.append("</section>")

    body = "\n".join(parts)
    doc = f"""<!doctype html><html lang="en"><head><meta charset="utf-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1"/>
<title>{e(title)}</title>
<style>
  :root{{--ink:{INK};--muted:{MUTED};--line:{LINE};--panel:{PANEL};--blue:{BLUE};}}
  *{{box-sizing:border-box;}}
  body{{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Inter,Arial,sans-serif;color:var(--ink);margin:0;background:#f1f5f9;}}
  .page{{max-width:900px;margin:0 auto;background:#fff;padding:56px 60px;}}
  .eyebrow{{color:var(--blue);font-weight:700;font-size:12px;letter-spacing:.16em;text-transform:uppercase;}}
  h1{{font-size:34px;line-height:1.1;margin:12px 0 8px;}}
  .sub{{color:var(--muted);font-size:16px;margin:0 0 14px;max-width:60ch;}}
  .meta{{color:var(--muted);font-size:12px;border-top:2px solid var(--ink);padding-top:14px;display:flex;gap:22px;}}
  section{{margin:34px 0;}}
  h2{{font-size:22px;border-bottom:1px solid var(--line);padding-bottom:8px;display:flex;gap:12px;align-items:baseline;}}
  h2 .num{{color:var(--blue);font-size:15px;font-weight:800;}}
  p{{font-size:14.5px;line-height:1.65;color:#334155;}}
  table{{width:100%;border-collapse:collapse;margin:12px 0;font-size:13px;box-shadow:0 1px 2px rgba(15,23,42,.04);}}
  th{{background:var(--blue);color:#fff;text-align:left;padding:9px 11px;font-weight:600;}}
  td{{padding:9px 11px;border-bottom:1px solid var(--line);vertical-align:top;color:#1e293b;}}
  tbody tr:nth-child(even){{background:var(--panel);}}
  .chart{{max-width:100%;height:auto;display:block;margin:10px auto;}}
  blockquote{{background:var(--panel);border-left:4px solid var(--blue);margin:14px 0;padding:14px 18px;font-size:17px;font-weight:700;}}
  blockquote cite{{display:block;font-size:12px;font-weight:400;color:var(--muted);margin-top:6px;font-style:normal;}}
  .src{{font-size:11.5px;color:var(--muted);}} .src a{{color:var(--blue);}}
  .footer{{margin-top:40px;border-top:1px solid var(--line);padding-top:14px;color:#94a3b8;font-size:11px;}}
  @media print{{body{{background:#fff;}}.page{{padding:0;max-width:none;}}}}
</style></head><body><div class="page">
  <div class="eyebrow">{e(studio)}</div>
  <h1>{e(title)}</h1>
  {f'<p class="sub">{e(subtitle)}</p>' if subtitle else ''}
  <div class="meta"><span><b>Aivora HC Advisory</b></span><span>{len(_sections(content))} sections</span><span>Confidential</span></div>
  {body}
  <div class="footer">{e(title)} &middot; Generated by Aivora HC</div>
</div></body></html>"""
    return doc.encode("utf-8")


# ── DOCX (python-docx) ────────────────────────────────────────────────────────

def build_docx(content: dict[str, Any], org_name: str = "") -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor, Inches

    title, subtitle, studio = _doc_meta(content)
    doc = Document()
    styles = doc.styles
    normal = styles["Normal"].font
    normal.name = "Calibri"
    normal.size = Pt(10.5)

    def _rgb(hex_):
        return RGBColor(int(hex_[1:3], 16), int(hex_[3:5], 16), int(hex_[5:7], 16))

    p = doc.add_paragraph()
    r = p.add_run(studio.upper())
    r.bold = True
    r.font.size = Pt(9)
    r.font.color.rgb = _rgb(BLUE)

    h = doc.add_paragraph()
    hr = h.add_run(title)
    hr.bold = True
    hr.font.size = Pt(22)
    hr.font.color.rgb = _rgb(INK)
    if subtitle:
        sp = doc.add_paragraph()
        sr = sp.add_run(subtitle)
        sr.font.size = Pt(11)
        sr.font.color.rgb = _rgb(MUTED)
    mp = doc.add_paragraph()
    mr = mp.add_run(f"Aivora HC Advisory   |   {len(_sections(content))} sections   |   Confidential")
    mr.font.size = Pt(8.5)
    mr.font.color.rgb = _rgb(MUTED)
    doc.add_paragraph().add_run().add_break()

    def _add_table(rows: list[list[str]]):
        header, *body_rows = rows
        table = doc.add_table(rows=1, cols=len(header))
        table.style = "Light Grid Accent 1"
        for j, htext in enumerate(header):
            c = table.rows[0].cells[j].paragraphs[0]
            rr = c.add_run(_clean(htext))
            rr.bold = True
        for br in body_rows:
            cells = table.add_row().cells
            for j, ctext in enumerate(br[:len(header)]):
                cells[j].paragraphs[0].add_run(_clean(ctext))
        doc.add_paragraph()

    for i, section in enumerate(_sections(content)):
        layout = str(section.get("layout") or "")
        data = section.get("data") if isinstance(section.get("data"), dict) else {}
        hp = doc.add_heading(level=2)
        hr2 = hp.add_run(f"{i+1:02d}   {_section_title(section, i)}")
        hr2.font.color.rgb = _rgb(INK)

        chart = report_charts.chart_png_for(layout, data)
        if chart:
            try:
                doc.add_picture(io.BytesIO(chart), width=Inches(6.2))
                doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
            except Exception:  # noqa: BLE001
                chart = None

        if layout in ("narrative_paragraph", "narrative", "prose", "text"):
            for para in _narrative_text(data).split("\n"):
                if para.strip():
                    doc.add_paragraph(_clean(para))
        elif layout in ("callout_quote", "quote", "callout", "note", "warning"):
            q, attr = _callout_text(data)
            qp = doc.add_paragraph()
            qr = qp.add_run(q)
            qr.bold = True
            qr.font.size = Pt(12.5)
            if attr:
                ap = doc.add_paragraph()
                ar = ap.add_run(attr)
                ar.font.size = Pt(9)
                ar.font.color.rgb = _rgb(MUTED)
        else:
            rows = _table_for(layout, data)
            if rows and len(rows) > 1:
                _add_table(rows)
            elif not chart:
                nb = _narrative_text(data)
                if nb:
                    doc.add_paragraph(nb)

        url = _footnote_url(section)
        if url:
            sp = doc.add_paragraph()
            sr = sp.add_run(f"Source: {url}")
            sr.font.size = Pt(8)
            sr.font.color.rgb = _rgb(MUTED)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PPTX (python-pptx) ────────────────────────────────────────────────────────

def build_pptx(content: dict[str, Any], org_name: str = "") -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    title, subtitle, studio = _doc_meta(content)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def _rgb(hex_):
        return RGBColor(int(hex_[1:3], 16), int(hex_[3:5], 16), int(hex_[5:7], 16))

    def _text(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=None):
        from pptx.enum.text import PP_ALIGN
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = _clean(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        if align:
            p.alignment = PP_ALIGN.LEFT
        return tb

    # Cover slide
    s0 = prs.slides.add_slide(blank)
    bg = s0.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb("#ffffff")
    _text(s0, Inches(0.8), Inches(2.4), SW - Inches(1.6), Inches(0.5), studio.upper(), size=14, bold=True, color=BLUE)
    _text(s0, Inches(0.8), Inches(2.9), SW - Inches(1.6), Inches(1.5), title, size=36, bold=True, color=INK)
    if subtitle:
        _text(s0, Inches(0.8), Inches(4.3), SW - Inches(1.6), Inches(1.2), subtitle, size=15, color=MUTED)
    _text(s0, Inches(0.8), Inches(6.5), SW - Inches(1.6), Inches(0.5),
          f"Aivora HC Advisory   |   {len(_sections(content))} sections   |   Confidential", size=11, color=MUTED)

    def _add_table(slide, rows, top):
        from pptx.util import Inches as In
        header, *body_rows = rows
        body_rows = body_rows[:10]  # keep slides readable
        nrows = len(body_rows) + 1
        ncols = len(header)
        gtbl = slide.shapes.add_table(nrows, ncols, In(0.6), top, SW - In(1.2), In(min(0.4 * nrows + 0.4, 5.5))).table
        for j, htext in enumerate(header):
            cell = gtbl.cell(0, j)
            cell.text = _clean(htext)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(BLUE)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = _rgb("#ffffff")
        for ri, br in enumerate(body_rows, start=1):
            for j in range(ncols):
                cell = gtbl.cell(ri, j)
                cell.text = _clean(br[j]) if j < len(br) else ""
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(9.5)
                para.font.color.rgb = _rgb(INK)

    for i, section in enumerate(_sections(content)):
        layout = str(section.get("layout") or "")
        data = section.get("data") if isinstance(section.get("data"), dict) else {}
        slide = prs.slides.add_slide(blank)
        _text(slide, Inches(0.6), Inches(0.35), SW - Inches(1.2), Inches(0.7),
              f"{i+1:02d}   {_section_title(section, i)}", size=22, bold=True, color=INK)

        chart = report_charts.chart_png_for(layout, data)
        placed = False
        if chart:
            try:
                slide.shapes.add_picture(io.BytesIO(chart), Inches(1.8), Inches(1.3), height=Inches(5.4))
                placed = True
            except Exception:  # noqa: BLE001
                chart = None
        if not placed:
            if layout in ("narrative_paragraph", "narrative", "prose", "text"):
                _text(slide, Inches(0.6), Inches(1.3), SW - Inches(1.2), Inches(5.5), _narrative_text(data), size=14, color="#334155")
            elif layout in ("callout_quote", "quote", "callout", "note", "warning"):
                q, attr = _callout_text(data)
                _text(slide, Inches(0.9), Inches(2.2), SW - Inches(1.8), Inches(2.5), q, size=22, bold=True, color=INK)
                if attr:
                    _text(slide, Inches(0.9), Inches(4.6), SW - Inches(1.8), Inches(0.8), attr, size=12, color=MUTED)
            else:
                rows = _table_for(layout, data)
                if rows and len(rows) > 1:
                    _add_table(slide, rows, Inches(1.3))
                else:
                    nb = _narrative_text(data)
                    if nb:
                        _text(slide, Inches(0.6), Inches(1.3), SW - Inches(1.2), Inches(5.5), nb, size=14, color="#334155")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
