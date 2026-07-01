"""
High-quality, format-differentiated exports.

Each generator is purpose-built for its medium:
- PPTX: cover slide + KPI slide + per-section content slides + closing slide,
        all with branded master colors and proper layout placeholders.
- PDF:  cover page + executive summary callout + recommendations grid with
        priority chips + per-section pages, sized for A4 print.
- DOCX: Word-styled report with heading hierarchy, an executive summary box,
        recommendation tables, and a page break per major section.
- HTML: standalone responsive web page mirroring the frontend infographic
        design, no external dependencies, suitable for client share.
"""
from __future__ import annotations

import html as _html
import io
from datetime import datetime, timezone
from typing import Any

SKIP_KEYS = {"skill_slug", "generated_at", "tool_name", "tool_input", "raw_text"}

# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

_SUMMARY_KEYS = ("executive_summary", "summary", "overview", "hc_vision", "strategy_vision", "industry_context")
_FINDINGS_KEYS = ("key_findings", "top_priorities", "quick_wins")
_RECS_KEYS = ("recommendations", "risks_and_mitigations", "priority_improvement_areas")
_ROADMAP_KEYS = ("implementation_plan", "maturity_roadmap", "implementation_roadmap", "12_month_roadmap", "90_day_plan")
_METRICS_KEYS = ("success_metrics", "metrics_and_targets", "benchmark_results")
_PILLAR_KEYS = ("strategic_pillars", "talent_lifecycle", "sourcing_strategy", "hiring_process", "investment_priorities")


def _source(content: dict[str, Any]) -> dict[str, Any]:
    """Return the renderable content, supporting old `tool_input` wrapper."""
    if isinstance(content.get("tool_input"), dict):
        return content["tool_input"]
    return {k: v for k, v in content.items() if k not in SKIP_KEYS and v}


def _title(content: dict[str, Any]) -> str:
    slug = content.get("skill_slug", "hc-advisory")
    return str(slug).replace("-", " ").title()


def _org_name(source: dict[str, Any], default: str) -> str:
    return str(source.get("org_name") or source.get("organisation") or default)


def _summary(source: dict[str, Any]) -> str | None:
    for k in _SUMMARY_KEYS:
        v = source.get(k)
        if isinstance(v, str) and v.strip():
            return v.strip()
    return None


def _key_section(source: dict[str, Any], keys: tuple[str, ...]) -> tuple[str | None, Any]:
    """Return the first (key, value) pair matching one of the given keys."""
    for k in keys:
        v = source.get(k)
        if v:
            return k, v
    return None, None


def _to_string_list(v: Any) -> list[str]:
    if isinstance(v, list):
        return [s for s in v if isinstance(s, str) and s.strip()]
    return []


def _to_dict_list(v: Any) -> list[dict[str, Any]]:
    if isinstance(v, list):
        return [d for d in v if isinstance(d, dict)]
    return []


def _pretty(k: str) -> str:
    return k.replace("_", " ").title()


def _safe(s: Any, limit: int = 4000) -> str:
    """Stringify safely, trimming long content."""
    if s is None:
        return ""
    if not isinstance(s, str):
        s = str(s)
    return s[:limit] + ("…" if len(s) > limit else "")


def _flatten_value(value: Any) -> str:
    """Used by content-slide bodies in PPTX where rich layout is not practical."""
    if isinstance(value, str):
        return value
    if isinstance(value, (int, float, bool)):
        return str(value)
    if isinstance(value, list):
        parts: list[str] = []
        for item in value:
            if isinstance(item, str):
                parts.append(f"• {item}")
            elif isinstance(item, dict):
                head = next(iter(item.values()), "")
                parts.append(f"• {head}")
        return "\n".join(parts)
    if isinstance(value, dict):
        return "\n".join(f"{_pretty(k)}: {v}" for k, v in value.items() if v)
    return str(value)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def generate_pptx(content: dict[str, Any], brand: dict[str, Any]) -> bytes:
    """Branded PPTX: cover + KPI strip slide + per-section slides + closing."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    def rgb(hx: str) -> RGBColor:
        h = hx.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    primary = brand.get("primary_color") or "#3b82f6"
    secondary = brand.get("secondary_color") or "#0ea5e9"
    font = brand.get("font") or "Calibri"

    source = _source(content)
    title = _title(content)
    org = _org_name(source, title)
    summary = _summary(source)
    findings_k, findings_v = _key_section(source, _FINDINGS_KEYS)
    recs_k, recs_v = _key_section(source, _RECS_KEYS)
    metrics_k, metrics_v = _key_section(source, _METRICS_KEYS)
    roadmap_k, roadmap_v = _key_section(source, _ROADMAP_KEYS)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    SH = prs.slide_height

    blank = prs.slide_layouts[6]

    def add_solid_rect(slide, x, y, w, h, color):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=None, align=None):
        from pptx.enum.text import PP_ALIGN
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        if color is not None:
            p.font.color.rgb = color
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        return box

    # ── 1. Cover slide ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_solid_rect(s, 0, 0, SW, Inches(0.4), rgb(primary))
    add_solid_rect(s, 0, SH - Inches(0.3), SW, Inches(0.3), rgb(secondary))
    add_text(s, Inches(0.8), Inches(1.0), SW - Inches(1.6), Inches(0.5), "AIVORA HC ADVISORY",
             size=11, bold=True, color=rgb(primary))
    add_text(s, Inches(0.8), Inches(1.7), SW - Inches(1.6), Inches(2.2), org,
             size=54, bold=True, color=RGBColor(0x0F, 0x17, 0x2A))
    add_text(s, Inches(0.8), Inches(3.7), SW - Inches(1.6), Inches(0.6), title,
             size=24, bold=False, color=rgb(primary))
    add_text(s, Inches(0.8), SH - Inches(1.0), SW - Inches(1.6), Inches(0.5),
             f"Generated {datetime.now(timezone.utc).strftime('%B %Y')}  ·  Confidential",
             size=11, color=RGBColor(0x64, 0x74, 0x8B))

    # ── 2. Executive summary slide ─────────────────────────────────────
    if summary:
        s = prs.slides.add_slide(blank)
        add_solid_rect(s, 0, 0, SW, Inches(0.6), rgb(primary))
        add_text(s, Inches(0.5), Inches(0.1), SW - Inches(1.0), Inches(0.4),
                 "EXECUTIVE SUMMARY", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        add_text(s, Inches(0.8), Inches(1.2), SW - Inches(1.6), SH - Inches(1.8),
                 _safe(summary, 1600), size=18, color=RGBColor(0x1E, 0x29, 0x3B))

    # ── 3. Metrics / KPI strip slide ───────────────────────────────────
    if metrics_v:
        s = prs.slides.add_slide(blank)
        add_solid_rect(s, 0, 0, SW, Inches(0.6), rgb(primary))
        add_text(s, Inches(0.5), Inches(0.1), SW, Inches(0.4),
                 (metrics_k or "metrics").replace("_", " ").upper(),
                 size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        metrics_list = _to_string_list(metrics_v)
        if not metrics_list:
            dicts = _to_dict_list(metrics_v)
            metrics_list = [
                f"{d.get('metric', '')}\n{d.get('target', d.get('baseline', d.get('organisation_value', '')))}"
                for d in dicts[:6]
            ]
        metrics_list = metrics_list[:6]
        if metrics_list:
            n = len(metrics_list)
            tile_w = (SW - Inches(1.2) - Inches(0.2) * (n - 1)) / n
            x = Inches(0.6)
            y = Inches(1.6)
            for m in metrics_list:
                tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, tile_w, Inches(2.2))
                tile.fill.solid()
                tile.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
                tile.line.color.rgb = rgb(primary)
                tile.line.width = Pt(0.75)
                add_text(s, x, y + Inches(0.8), tile_w, Inches(1.0), m,
                         size=14, bold=True, color=RGBColor(0x1E, 0x29, 0x3B), align="center")
                x += tile_w + Inches(0.2)

    # ── 4. Findings ────────────────────────────────────────────────────
    if findings_v:
        s = prs.slides.add_slide(blank)
        add_solid_rect(s, 0, 0, SW, Inches(0.6), rgb(primary))
        add_text(s, Inches(0.5), Inches(0.1), SW, Inches(0.4),
                 (findings_k or "key findings").replace("_", " ").upper(),
                 size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        finds = _to_string_list(findings_v)[:6]
        y = Inches(1.2)
        for i, f in enumerate(finds, start=1):
            num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.6), Inches(0.6))
            num.fill.solid()
            num.fill.fore_color.rgb = rgb(primary)
            num.line.fill.background()
            add_text(s, Inches(0.7), y, Inches(0.6), Inches(0.6), str(i),
                     size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align="center")
            add_text(s, Inches(1.5), y, SW - Inches(2.0), Inches(0.8), f,
                     size=14, color=RGBColor(0x1E, 0x29, 0x3B))
            y += Inches(0.9)

    # ── 5. Recommendations ─────────────────────────────────────────────
    if recs_v:
        recs = _to_dict_list(recs_v)[:4]
        s = prs.slides.add_slide(blank)
        add_solid_rect(s, 0, 0, SW, Inches(0.6), rgb(primary))
        add_text(s, Inches(0.5), Inches(0.1), SW, Inches(0.4),
                 (recs_k or "recommendations").replace("_", " ").upper(),
                 size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        # 2-col grid of recommendation cards
        col_w = (SW - Inches(1.6)) / 2
        for i, r in enumerate(recs):
            col = i % 2
            row = i // 2
            x = Inches(0.5) + col * (col_w + Inches(0.3))
            y = Inches(1.2) + row * Inches(2.7)
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, Inches(2.5))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            card.line.width = Pt(0.75)
            t = str(r.get("recommendation") or r.get("title") or r.get("area") or r.get("risk") or "")
            add_text(s, x + Inches(0.2), y + Inches(0.2), col_w - Inches(0.4), Inches(0.6),
                     t, size=14, bold=True, color=RGBColor(0x1E, 0x29, 0x3B))
            rat = str(r.get("rationale") or r.get("mitigation") or r.get("description") or "")
            if rat:
                add_text(s, x + Inches(0.2), y + Inches(0.9), col_w - Inches(0.4), Inches(1.4),
                         _safe(rat, 280), size=11, color=RGBColor(0x47, 0x55, 0x69))
            pri = str(r.get("priority", "")).lower()
            if pri:
                pri_color = (
                    RGBColor(0xDC, 0x26, 0x26) if "high" in pri
                    else RGBColor(0x10, 0xB9, 0x81) if "low" in pri
                    else RGBColor(0xF5, 0x9E, 0x0B)
                )
                add_text(s, x + Inches(0.2), y + Inches(2.0), col_w - Inches(0.4), Inches(0.3),
                         pri.upper(), size=10, bold=True, color=pri_color)

    # ── 6. Roadmap ─────────────────────────────────────────────────────
    if roadmap_v:
        rows = _to_dict_list(roadmap_v)[:5]
        s = prs.slides.add_slide(blank)
        add_solid_rect(s, 0, 0, SW, Inches(0.6), rgb(primary))
        add_text(s, Inches(0.5), Inches(0.1), SW, Inches(0.4),
                 (roadmap_k or "roadmap").replace("_", " ").upper(),
                 size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        # Horizontal pipeline of phase tiles
        if rows:
            n = len(rows)
            tile_w = (SW - Inches(1.0) - Inches(0.15) * (n - 1)) / n
            x = Inches(0.5)
            y = Inches(2.5)
            for i, r in enumerate(rows):
                phase = str(r.get("phase") or r.get("week_range") or f"Phase {i + 1}")
                tl = str(r.get("timeline") or "")
                tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, tile_w, Inches(2.5))
                tile.fill.solid()
                tile.fill.fore_color.rgb = RGBColor(0x10, 0xB9, 0x81) if i % 2 == 0 else rgb(secondary)
                tile.line.fill.background()
                add_text(s, x, y + Inches(0.3), tile_w, Inches(0.6), phase,
                         size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align="center")
                if tl:
                    add_text(s, x, y + Inches(1.0), tile_w, Inches(0.5), tl,
                             size=11, color=RGBColor(0xFF, 0xFF, 0xFF), align="center")
                actions = _to_string_list(r.get("actions") or r.get("milestones") or r.get("focus_areas")) [:2]
                if actions:
                    add_text(s, x + Inches(0.2), y + Inches(1.5), tile_w - Inches(0.4), Inches(0.9),
                             "  ·  ".join(actions), size=10,
                             color=RGBColor(0xFF, 0xFF, 0xFF), align="center")
                x += tile_w + Inches(0.15)

    # ── 7. Remaining sections — one slide each ─────────────────────────
    used = {findings_k, recs_k, metrics_k, roadmap_k} | set(_SUMMARY_KEYS)
    for key, value in source.items():
        if key in used or not value:
            continue
        if key in (findings_k, recs_k, metrics_k, roadmap_k):
            continue
        body = _flatten_value(value)
        if not body.strip():
            continue
        s = prs.slides.add_slide(blank)
        add_solid_rect(s, 0, 0, SW, Inches(0.6), rgb(primary))
        add_text(s, Inches(0.5), Inches(0.1), SW, Inches(0.4),
                 key.replace("_", " ").upper(), size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
        add_text(s, Inches(0.7), Inches(1.0), SW - Inches(1.4), SH - Inches(1.4),
                 _safe(body, 2200), size=13, color=RGBColor(0x1E, 0x29, 0x3B))

    # ── 8. Closing slide ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_solid_rect(s, 0, 0, SW, SH, rgb(primary))
    add_text(s, Inches(0.8), Inches(2.7), SW - Inches(1.6), Inches(1.0),
             "Thank you.", size=60, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, Inches(0.8), Inches(4.0), SW - Inches(1.6), Inches(0.6),
             "Aivora HC Advisory  ·  Confidential",
             size=14, color=RGBColor(0xE2, 0xE8, 0xF0))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def generate_pdf(content: dict[str, Any], brand: dict[str, Any]) -> bytes:
    """Print-quality PDF: cover page, executive callout, recommendation cards, sections."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        HRFlowable, KeepTogether, PageBreak, Paragraph, SimpleDocTemplate,
        Spacer, Table, TableStyle,
    )

    primary_hex = brand.get("primary_color") or "#3b82f6"
    secondary_hex = brand.get("secondary_color") or "#0ea5e9"
    primary = colors.HexColor(primary_hex if primary_hex.startswith("#") else f"#{primary_hex}")
    secondary = colors.HexColor(secondary_hex if secondary_hex.startswith("#") else f"#{secondary_hex}")

    source = _source(content)
    title = _title(content)
    org = _org_name(source, title)
    summary = _summary(source)
    findings_k, findings_v = _key_section(source, _FINDINGS_KEYS)
    recs_k, recs_v = _key_section(source, _RECS_KEYS)
    metrics_k, metrics_v = _key_section(source, _METRICS_KEYS)
    roadmap_k, roadmap_v = _key_section(source, _ROADMAP_KEYS)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=2 * cm, rightMargin=2 * cm,
                            topMargin=2 * cm, bottomMargin=2 * cm)

    styles = getSampleStyleSheet()
    h_org = ParagraphStyle("HOrg", parent=styles["Heading1"], fontSize=36,
                           textColor=colors.HexColor("#0F172A"), leading=42, spaceAfter=4)
    h_title = ParagraphStyle("HT", parent=styles["Heading2"], fontSize=18,
                             textColor=primary, leading=22, spaceAfter=20)
    h_section = ParagraphStyle("HS", parent=styles["Heading2"], fontSize=14,
                               textColor=primary, leading=18, spaceBefore=16, spaceAfter=8)
    label = ParagraphStyle("Label", parent=styles["Normal"], fontSize=9,
                           textColor=primary, leading=12, spaceAfter=4)
    body = ParagraphStyle("Body", parent=styles["Normal"], fontSize=10,
                          leading=15, textColor=colors.HexColor("#1E293B"), spaceAfter=6)
    body_small = ParagraphStyle("BodyS", parent=body, fontSize=9, leading=13)
    bullet = ParagraphStyle("Bullet", parent=body, leftIndent=14, bulletText="▪", bulletIndent=4)
    meta = ParagraphStyle("Meta", parent=styles["Normal"], fontSize=9,
                          textColor=colors.HexColor("#64748B"))
    summary_style = ParagraphStyle("Sum", parent=body, fontSize=11, leading=17,
                                   textColor=colors.HexColor("#1E293B"), backColor=colors.HexColor("#F1F5F9"),
                                   borderColor=primary, borderWidth=0, borderPadding=10,
                                   leftIndent=0, rightIndent=0, spaceBefore=4, spaceAfter=8)

    story: list[Any] = []

    # ── Cover ─────────────────────────────────────────────────────────
    story.append(Paragraph("AIVORA HC ADVISORY", label))
    story.append(HRFlowable(width="20%", thickness=3, color=primary, spaceBefore=2, spaceAfter=14))
    story.append(Paragraph(_safe(org, 120), h_org))
    story.append(Paragraph(_safe(title, 120), h_title))
    story.append(Paragraph(
        f"Generated {datetime.now(timezone.utc).strftime('%B %Y')}  ·  Confidential", meta))
    story.append(Spacer(1, 0.4 * cm))
    story.append(HRFlowable(width="100%", thickness=2, color=primary))

    # ── Executive summary callout ─────────────────────────────────────
    if summary:
        story.append(Paragraph("EXECUTIVE SUMMARY", label))
        story.append(Paragraph(_safe(summary, 1800), summary_style))

    # ── KPI strip ─────────────────────────────────────────────────────
    if metrics_v:
        story.append(Paragraph((metrics_k or "metrics").replace("_", " ").upper(), label))
        story.append(Paragraph(_pretty(metrics_k or "Metrics"), h_section))
        cells = _to_string_list(metrics_v)
        if not cells:
            cells = [
                f"<b>{d.get('metric', '')}</b><br/>{d.get('target', d.get('baseline', ''))}"
                for d in _to_dict_list(metrics_v)[:6]
            ]
        cells = cells[:6]
        if cells:
            cols = min(3, len(cells))
            rows = [cells[i:i + cols] for i in range(0, len(cells), cols)]
            for row in rows:
                while len(row) < cols:
                    row.append("")
            tbl_data = [[Paragraph(c, body_small) for c in r] for r in rows]
            tbl = Table(tbl_data, colWidths=[(A4[0] - 4 * cm) / cols] * cols)
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F1F5F9")),
                ("BOX", (0, 0), (-1, -1), 0.5, primary),
                ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#E2E8F0")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("TOPPADDING", (0, 0), (-1, -1), 12),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 12),
            ]))
            story.append(KeepTogether([tbl]))

    # ── Findings ──────────────────────────────────────────────────────
    if findings_v:
        story.append(Paragraph((findings_k or "findings").replace("_", " ").upper(), label))
        story.append(Paragraph(_pretty(findings_k or "Key findings"), h_section))
        for i, f in enumerate(_to_string_list(findings_v)[:8], start=1):
            data = [[Paragraph(f"<b><font color='{primary_hex}'>{i}</font></b>", body),
                     Paragraph(_safe(f, 300), body)]]
            tbl = Table(data, colWidths=[0.8 * cm, A4[0] - 4 * cm - 0.8 * cm])
            tbl.setStyle(TableStyle([
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("BACKGROUND", (0, 0), (0, 0), colors.HexColor("#FEF3C7")),
                ("BOX", (0, 0), (-1, -1), 0.25, colors.HexColor("#E2E8F0")),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 0.15 * cm))

    # ── Recommendations ──────────────────────────────────────────────
    if recs_v:
        story.append(Spacer(1, 0.3 * cm))
        story.append(Paragraph((recs_k or "recommendations").replace("_", " ").upper(), label))
        story.append(Paragraph(_pretty(recs_k or "Recommendations"), h_section))
        for r in _to_dict_list(recs_v)[:6]:
            t = str(r.get("recommendation") or r.get("title") or r.get("area") or r.get("risk") or "")
            rat = str(r.get("rationale") or r.get("mitigation") or r.get("description") or "")
            pri = str(r.get("priority", "")).lower()
            pri_hex = "#DC2626" if "high" in pri else "#10B981" if "low" in pri else "#F59E0B"
            chip_html = (f'<para><font color="{pri_hex}"><b>{pri.upper()}</b></font></para>'
                         if pri else "")
            data = [
                [Paragraph(f"<b>{_safe(t, 200)}</b>", body),
                 Paragraph(chip_html, body_small)],
                [Paragraph(_safe(rat, 400), body), ""],
            ]
            tbl = Table(data, colWidths=[(A4[0] - 4 * cm) * 0.78, (A4[0] - 4 * cm) * 0.22])
            tbl.setStyle(TableStyle([
                ("SPAN", (0, 1), (1, 1)),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ALIGN", (1, 0), (1, 0), "RIGHT"),
                ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#E2E8F0")),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F8FAFC")),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                ("LEFTPADDING", (0, 0), (-1, -1), 10),
                ("RIGHTPADDING", (0, 0), (-1, -1), 10),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 0.2 * cm))

    # ── Roadmap ──────────────────────────────────────────────────────
    if roadmap_v:
        story.append(PageBreak())
        story.append(Paragraph((roadmap_k or "roadmap").replace("_", " ").upper(), label))
        story.append(Paragraph(_pretty(roadmap_k or "Roadmap"), h_section))
        for i, r in enumerate(_to_dict_list(roadmap_v)[:8], start=1):
            phase = str(r.get("phase") or r.get("week_range") or f"Phase {i}")
            tl = str(r.get("timeline") or "")
            head_html = f"<b>{_safe(phase, 80)}</b>"
            if tl:
                head_html += f"  ·  <font color='{primary_hex}'>{_safe(tl, 80)}</font>"
            story.append(Paragraph(head_html, body))
            actions = _to_string_list(r.get("actions") or r.get("milestones") or r.get("focus_areas"))
            for a in actions[:4]:
                story.append(Paragraph(_safe(a, 240), bullet))
            outcome = str(r.get("expected_outcome") or "")
            if outcome:
                story.append(Paragraph(f"<i>Outcome: {_safe(outcome, 240)}</i>", body_small))
            story.append(Spacer(1, 0.18 * cm))

    # ── Remaining sections ──────────────────────────────────────────
    used = {findings_k, recs_k, metrics_k, roadmap_k} | set(_SUMMARY_KEYS)
    rest = [(k, v) for k, v in source.items() if k not in used and v]
    if rest:
        story.append(PageBreak())
        for key, value in rest:
            story.append(Paragraph(key.replace("_", " ").upper(), label))
            story.append(Paragraph(_pretty(key), h_section))
            if isinstance(value, str):
                story.append(Paragraph(_safe(value, 3000), body))
            elif isinstance(value, list):
                for item in value:
                    if isinstance(item, str):
                        story.append(Paragraph(_safe(item, 800), bullet))
                    elif isinstance(item, dict):
                        line = ", ".join(f"<b>{_pretty(k)}:</b> {v}" for k, v in list(item.items())[:3] if v)
                        story.append(Paragraph(_safe(line, 800), bullet))
            elif isinstance(value, dict):
                for k, v in value.items():
                    if v:
                        story.append(Paragraph(f"<b>{_pretty(k)}:</b> {_safe(v, 600)}", body))

    # ── Footer ──────────────────────────────────────────────────────
    story.append(Spacer(1, 0.6 * cm))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#CBD5E1")))
    story.append(Paragraph(
        f"Aivora HC Platform  ·  Confidential  ·  {datetime.now(timezone.utc).year}", meta))

    # quiet linter: keep secondary referenced
    _ = secondary
    doc.build(story)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def generate_docx(content: dict[str, Any], brand: dict[str, Any]) -> bytes:
    """Word document with executive summary box, recommendation tables, section breaks."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import Pt, RGBColor, Cm

    primary_hex = (brand.get("primary_color") or "#3b82f6").lstrip("#")
    secondary_hex = (brand.get("secondary_color") or "#0ea5e9").lstrip("#")
    primary = RGBColor(int(primary_hex[0:2], 16), int(primary_hex[2:4], 16), int(primary_hex[4:6], 16))
    secondary = RGBColor(int(secondary_hex[0:2], 16), int(secondary_hex[2:4], 16), int(secondary_hex[4:6], 16))
    font = brand.get("font") or "Calibri"

    source = _source(content)
    title = _title(content)
    org = _org_name(source, title)
    summary = _summary(source)
    findings_k, findings_v = _key_section(source, _FINDINGS_KEYS)
    recs_k, recs_v = _key_section(source, _RECS_KEYS)
    metrics_k, metrics_v = _key_section(source, _METRICS_KEYS)
    roadmap_k, roadmap_v = _key_section(source, _ROADMAP_KEYS)

    doc = Document()

    # Tweak default font
    style = doc.styles["Normal"]
    style.font.name = font
    style.font.size = Pt(11)

    def add_para(text: str, *, size=11, bold=False, color: RGBColor | None = None,
                 align=None, space_after=4):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        if align == "center":
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(space_after)
        return p

    def add_heading(text: str, *, level: int = 1):
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.name = font
            run.font.color.rgb = primary
            if level == 0:
                run.font.size = Pt(28)
            elif level == 1:
                run.font.size = Pt(16)
            else:
                run.font.size = Pt(13)
        return h

    def set_cell_bg(cell, hex_color: str):
        tc_pr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color.lstrip("#"))
        tc_pr.append(shd)

    # ── Cover ─────────────────────────────────────────────────────────
    add_para("AIVORA HC ADVISORY", size=10, bold=True, color=primary, space_after=4)
    add_heading(org, level=0)
    add_para(title, size=14, color=primary, space_after=12)
    add_para(f"Generated {datetime.now(timezone.utc).strftime('%B %Y')}  ·  Confidential",
             size=9, color=RGBColor(0x64, 0x74, 0x8B), space_after=12)

    # ── Executive summary as a shaded one-cell table ─────────────────
    if summary:
        add_heading("Executive summary", level=1)
        tbl = doc.add_table(rows=1, cols=1)
        cell = tbl.rows[0].cells[0]
        set_cell_bg(cell, "F1F5F9")
        cell.text = ""
        para = cell.paragraphs[0]
        run = para.add_run(_safe(summary, 2400))
        run.font.name = font
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        doc.add_paragraph()

    # ── Findings — numbered list ────────────────────────────────────
    if findings_v:
        add_heading(_pretty(findings_k or "Key findings"), level=1)
        for i, f in enumerate(_to_string_list(findings_v)[:10], start=1):
            p = doc.add_paragraph()
            run_no = p.add_run(f"{i}.  ")
            run_no.bold = True
            run_no.font.color.rgb = primary
            run_no.font.name = font
            run_text = p.add_run(_safe(f, 500))
            run_text.font.name = font
            run_text.font.size = Pt(11)

    # ── Recommendations table ───────────────────────────────────────
    if recs_v:
        add_heading(_pretty(recs_k or "Recommendations"), level=1)
        recs = _to_dict_list(recs_v)[:10]
        tbl = doc.add_table(rows=1, cols=3)
        tbl.style = "Light Grid Accent 1"
        hdr = tbl.rows[0].cells
        for i, h in enumerate(["Recommendation", "Rationale", "Priority"]):
            hdr[i].text = h
            for para in hdr[i].paragraphs:
                for run in para.runs:
                    run.font.name = font
                    run.font.bold = True
                    run.font.size = Pt(10)
        for r in recs:
            row = tbl.add_row().cells
            row[0].text = _safe(str(r.get("recommendation") or r.get("title") or r.get("area") or r.get("risk") or ""), 300)
            row[1].text = _safe(str(r.get("rationale") or r.get("mitigation") or r.get("description") or ""), 400)
            row[2].text = _safe(str(r.get("priority", "")).upper(), 30)

    # ── Metrics — KPI grid (small table) ────────────────────────────
    if metrics_v:
        add_heading(_pretty(metrics_k or "Metrics"), level=1)
        cells = _to_string_list(metrics_v)
        if not cells:
            cells = [
                f"{d.get('metric', '')}: {d.get('target', d.get('baseline', ''))}"
                for d in _to_dict_list(metrics_v)
            ]
        for c in cells[:10]:
            p = doc.add_paragraph(c)
            p.style = doc.styles["List Bullet"]

    # ── Roadmap ─────────────────────────────────────────────────────
    if roadmap_v:
        add_heading(_pretty(roadmap_k or "Roadmap"), level=1)
        for i, r in enumerate(_to_dict_list(roadmap_v)[:10], start=1):
            phase = str(r.get("phase") or r.get("week_range") or f"Phase {i}")
            tl = str(r.get("timeline") or "")
            p = doc.add_paragraph()
            run = p.add_run(f"{phase}")
            run.font.name = font
            run.bold = True
            run.font.color.rgb = primary
            run.font.size = Pt(12)
            if tl:
                run2 = p.add_run(f"   ·  {tl}")
                run2.font.color.rgb = secondary
                run2.font.size = Pt(10)
            actions = _to_string_list(r.get("actions") or r.get("milestones") or r.get("focus_areas"))
            for a in actions[:5]:
                bp = doc.add_paragraph(a)
                bp.style = doc.styles["List Bullet"]
            outcome = str(r.get("expected_outcome") or "")
            if outcome:
                op = doc.add_paragraph()
                run = op.add_run(f"Outcome: {outcome}")
                run.italic = True
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0x47, 0x55, 0x69)

    # ── Remaining sections ────────────────────────────────────────
    used = {findings_k, recs_k, metrics_k, roadmap_k} | set(_SUMMARY_KEYS)
    rest = [(k, v) for k, v in source.items() if k not in used and v]
    for key, value in rest:
        add_heading(_pretty(key), level=1)
        if isinstance(value, str):
            for line in value.split("\n"):
                if line.strip():
                    doc.add_paragraph(line.strip())
        elif isinstance(value, list):
            for item in value:
                if isinstance(item, str):
                    p = doc.add_paragraph(item)
                    p.style = doc.styles["List Bullet"]
                elif isinstance(item, dict):
                    p = doc.add_paragraph()
                    for k, v in list(item.items())[:4]:
                        if v:
                            run = p.add_run(f"{_pretty(k)}: ")
                            run.bold = True
                            p.add_run(f"{v}  ")
        elif isinstance(value, dict):
            for k, v in value.items():
                if v:
                    p = doc.add_paragraph()
                    run = p.add_run(f"{_pretty(k)}: ")
                    run.bold = True
                    p.add_run(_safe(str(v), 600))

    # ── Footer ───────────────────────────────────────────────────
    doc.add_paragraph()
    foot = doc.add_paragraph()
    fr = foot.add_run(f"Aivora HC Platform  ·  Confidential  ·  {datetime.now(timezone.utc).year}")
    fr.font.size = Pt(9)
    fr.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

    # silence linter
    _ = Cm

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def generate_xlsx(content: dict[str, Any], brand: dict[str, Any]) -> bytes:
    """Workbook with one sheet per section.

    KPI metrics become a tabular sheet, list-of-dicts become tables with one row
    per item, list-of-strings become a single column, plain strings/dicts fall
    back to a key/value layout. Anything we cannot fit cleanly is dropped onto
    a "Notes" sheet as text so no section is silently lost.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    primary_hex = (brand.get("primary_color") or "#3b82f6").lstrip("#")
    font_name = brand.get("font") or "Calibri"

    source = _source(content)
    title = _title(content)
    org = _org_name(source, title)
    summary = _summary(source)

    wb = Workbook()
    # Remove default sheet — we add our own.
    default = wb.active
    wb.remove(default)

    header_fill = PatternFill("solid", fgColor=primary_hex)
    header_font = Font(name=font_name, bold=True, color="FFFFFF", size=11)
    title_font = Font(name=font_name, bold=True, size=18, color="0F172A")
    sub_font = Font(name=font_name, size=11, color="64748B")
    wrap = Alignment(wrap_text=True, vertical="top")

    used_names: set[str] = set()

    def _sheet_name(name: str) -> str:
        # Excel sheet names: max 31 chars, no : \ / ? * [ ]
        clean = "".join(c for c in name if c not in ":\\/?*[]")[:31].strip() or "Sheet"
        base = clean
        i = 2
        while clean.lower() in used_names:
            suffix = f" {i}"
            clean = (base[: 31 - len(suffix)] + suffix).strip()
            i += 1
        used_names.add(clean.lower())
        return clean

    def _autosize(ws, max_width: int = 60) -> None:
        for col_cells in ws.columns:
            letter = get_column_letter(col_cells[0].column)
            longest = 0
            for cell in col_cells:
                v = cell.value
                if v is None:
                    continue
                # only sample first line; openpyxl can choke on huge strings
                line_len = max((len(s) for s in str(v).split("\n")), default=0)
                if line_len > longest:
                    longest = line_len
            ws.column_dimensions[letter].width = min(max(12, longest + 2), max_width)

    def _header_row(ws, headers: list[str]) -> None:
        for i, h in enumerate(headers, start=1):
            c = ws.cell(row=1, column=i, value=h)
            c.fill = header_fill
            c.font = header_font
            c.alignment = wrap

    # ── Overview ────────────────────────────────────────────────────
    ws = wb.create_sheet(_sheet_name("Overview"))
    ws["A1"] = org
    ws["A1"].font = title_font
    ws["A2"] = title
    ws["A2"].font = Font(name=font_name, size=13, color=primary_hex, bold=True)
    ws["A3"] = f"Generated {datetime.now(timezone.utc).strftime('%B %Y')}  ·  Confidential"
    ws["A3"].font = sub_font
    if summary:
        ws["A5"] = "Executive summary"
        ws["A5"].font = Font(name=font_name, bold=True, size=12, color=primary_hex)
        ws["A6"] = _safe(summary, 4000)
        ws["A6"].alignment = wrap
        ws.row_dimensions[6].height = 120
    ws.column_dimensions["A"].width = 100

    def _add_dict_list_sheet(name: str, rows: list[dict[str, Any]]) -> None:
        # Union of keys, ordered by first occurrence.
        headers: list[str] = []
        for r in rows:
            for k in r.keys():
                if k not in headers:
                    headers.append(k)
        if not headers:
            return
        sh = wb.create_sheet(_sheet_name(name))
        _header_row(sh, [_pretty(h) for h in headers])
        for ri, r in enumerate(rows, start=2):
            for ci, h in enumerate(headers, start=1):
                v = r.get(h)
                if isinstance(v, (list, dict)):
                    v = _flatten_value(v)
                cell = sh.cell(row=ri, column=ci, value=_safe(v, 2000) if v is not None else "")
                cell.alignment = wrap
        _autosize(sh)

    def _add_string_list_sheet(name: str, items: list[str]) -> None:
        sh = wb.create_sheet(_sheet_name(name))
        _header_row(sh, [_pretty(name)])
        for ri, s in enumerate(items, start=2):
            cell = sh.cell(row=ri, column=1, value=_safe(s, 2000))
            cell.alignment = wrap
        _autosize(sh, max_width=100)

    def _add_kv_sheet(name: str, mapping: dict[str, Any]) -> None:
        sh = wb.create_sheet(_sheet_name(name))
        _header_row(sh, ["Field", "Value"])
        ri = 2
        for k, v in mapping.items():
            if not v:
                continue
            sh.cell(row=ri, column=1, value=_pretty(str(k)))
            cell = sh.cell(row=ri, column=2, value=_safe(_flatten_value(v), 4000))
            cell.alignment = wrap
            ri += 1
        _autosize(sh, max_width=80)

    def _add_text_sheet(name: str, text: str) -> None:
        sh = wb.create_sheet(_sheet_name(name))
        _header_row(sh, [_pretty(name)])
        sh.cell(row=2, column=1, value=_safe(text, 16000)).alignment = wrap
        sh.column_dimensions["A"].width = 100

    # ── One sheet per top-level section ─────────────────────────────
    # Skip the summary keys (rendered on Overview) and ignore empty values.
    rendered_any = False
    for key, value in source.items():
        if not value:
            continue
        if key in _SUMMARY_KEYS:
            continue
        try:
            if isinstance(value, list):
                dicts = _to_dict_list(value)
                strings = _to_string_list(value)
                if dicts:
                    _add_dict_list_sheet(key, dicts)
                elif strings:
                    _add_string_list_sheet(key, strings)
                else:
                    # mixed / unknown — fall back to flattened text
                    _add_text_sheet(key, _flatten_value(value))
            elif isinstance(value, dict):
                _add_kv_sheet(key, value)
            elif isinstance(value, str):
                _add_text_sheet(key, value)
            else:
                _add_text_sheet(key, str(value))
            rendered_any = True
        except Exception as exc:  # noqa: BLE001 — never let one bad section drop the whole file
            _add_text_sheet(f"{key} (raw)", f"{_flatten_value(value)}\n\n[render error: {exc}]")
            rendered_any = True

    if not rendered_any and not summary:
        # Defensive: ensure the workbook has at least one informative sheet.
        sh = wb.create_sheet(_sheet_name("Content"))
        sh["A1"] = "No structured content available."
        sh["A1"].font = sub_font

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def generate_html(content: dict[str, Any], brand: dict[str, Any]) -> bytes:
    """Standalone responsive HTML page that mirrors the frontend infographic."""
    primary = brand.get("primary_color") or "#3b82f6"
    secondary = brand.get("secondary_color") or "#0ea5e9"
    accent_text = "#0f172a"

    source = _source(content)
    title = _title(content)
    org = _org_name(source, title)
    summary = _summary(source)
    findings_k, findings_v = _key_section(source, _FINDINGS_KEYS)
    recs_k, recs_v = _key_section(source, _RECS_KEYS)
    metrics_k, metrics_v = _key_section(source, _METRICS_KEYS)
    roadmap_k, roadmap_v = _key_section(source, _ROADMAP_KEYS)

    def e(s: Any) -> str:
        return _html.escape(str(s) if s is not None else "")

    parts: list[str] = []

    parts.append(f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8" />
<meta name="viewport" content="width=device-width, initial-scale=1" />
<title>{e(org)} - {e(title)}</title>
<style>
:root {{
  --primary: {primary};
  --secondary: {secondary};
  --ink: {accent_text};
}}
* {{ box-sizing: border-box; }}
body {{ margin: 0; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif; color: var(--ink); background: #f8fafc; line-height: 1.6; }}
.wrap {{ max-width: 980px; margin: 0 auto; padding: 32px 24px 80px; }}
.cover {{ padding: 60px 0 36px; border-bottom: 1px solid #e2e8f0; }}
.cover .eyebrow {{ font-size: 11px; font-weight: 700; letter-spacing: 0.18em; color: var(--primary); text-transform: uppercase; }}
.cover h1 {{ font-size: 44px; line-height: 1.1; margin: 8px 0 6px; letter-spacing: -0.02em; }}
.cover .sub {{ font-size: 20px; color: var(--primary); font-weight: 500; }}
.cover .meta {{ color: #64748b; font-size: 14px; margin-top: 18px; }}
.bar {{ height: 4px; background: linear-gradient(90deg, var(--primary), var(--secondary)); border-radius: 999px; margin-top: 14px; }}
section {{ margin: 36px 0; }}
.label {{ font-size: 11px; font-weight: 700; letter-spacing: 0.18em; color: var(--primary); text-transform: uppercase; margin-bottom: 12px; display: inline-flex; align-items: center; gap: 8px; }}
.label::before {{ content: ""; width: 22px; height: 2px; background: var(--primary); border-radius: 999px; }}
.callout {{ background: #f1f5f9; border-left: 4px solid var(--primary); padding: 20px 22px; border-radius: 12px; font-size: 16px; }}
.kpi-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr)); gap: 12px; }}
.kpi {{ background: #fff; border: 1px solid #e2e8f0; border-radius: 16px; padding: 18px; text-align: center; }}
.kpi h4 {{ margin: 0 0 6px; font-size: 11px; font-weight: 700; color: var(--primary); text-transform: uppercase; letter-spacing: 0.12em; }}
.kpi .v {{ font-size: 15px; color: var(--ink); }}
.find-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(260px, 1fr)); gap: 12px; }}
.find {{ background: #fff; border: 1px solid #e2e8f0; border-radius: 16px; padding: 18px; display: grid; grid-template-columns: 36px 1fr; gap: 14px; align-items: start; }}
.find .num {{ width: 32px; height: 32px; border-radius: 10px; background: #fef3c7; color: #b45309; font-weight: 800; display: grid; place-items: center; }}
.rec {{ background: #fff; border: 1px solid #e2e8f0; border-radius: 16px; padding: 18px 20px; margin-bottom: 12px; }}
.rec .top {{ display: flex; justify-content: space-between; gap: 12px; align-items: start; margin-bottom: 8px; }}
.rec h3 {{ margin: 0; font-size: 16px; }}
.chip {{ font-size: 10px; font-weight: 800; letter-spacing: 0.12em; padding: 4px 10px; border-radius: 999px; text-transform: uppercase; }}
.chip-high {{ background: #fee2e2; color: #b91c1c; }}
.chip-med {{ background: #fef3c7; color: #b45309; }}
.chip-low {{ background: #dcfce7; color: #166534; }}
.rec p {{ margin: 4px 0 0; color: #475569; font-size: 14px; }}
.rec .chips {{ display: flex; flex-wrap: wrap; gap: 6px; margin-top: 10px; }}
.rec .chips span {{ background: #eff6ff; color: #1d4ed8; font-size: 12px; padding: 4px 10px; border-radius: 999px; }}
.timeline {{ position: relative; padding-left: 22px; }}
.timeline::before {{ content: ""; position: absolute; left: 8px; top: 6px; bottom: 6px; width: 2px; background: linear-gradient(180deg, var(--primary), transparent); }}
.phase {{ position: relative; background: #fff; border: 1px solid #e2e8f0; border-radius: 16px; padding: 16px 18px; margin-bottom: 12px; }}
.phase::before {{ content: ""; position: absolute; left: -18px; top: 20px; width: 12px; height: 12px; border-radius: 999px; background: var(--primary); box-shadow: 0 0 0 4px rgba(59,130,246,0.15); }}
.phase h3 {{ margin: 0 0 4px; font-size: 15px; }}
.phase .tl {{ font-size: 12px; color: var(--primary); margin-bottom: 8px; }}
.phase ul {{ margin: 6px 0 0; padding-left: 18px; color: #334155; font-size: 14px; }}
.section-card {{ background: #fff; border: 1px solid #e2e8f0; border-radius: 16px; padding: 20px; margin-bottom: 12px; }}
.section-card h2 {{ font-size: 16px; margin: 0 0 8px; color: var(--primary); }}
footer {{ margin-top: 60px; color: #94a3b8; font-size: 12px; border-top: 1px solid #e2e8f0; padding-top: 16px; text-align: center; }}
@media print {{ body {{ background: #fff; }} .wrap {{ padding: 24px; }} }}
</style></head><body><div class="wrap">""")

    # Cover
    parts.append(f"""<div class="cover">
      <div class="eyebrow">Aivora HC Advisory</div>
      <h1>{e(org)}</h1>
      <div class="sub">{e(title)}</div>
      <div class="bar"></div>
      <div class="meta">Generated {datetime.now(timezone.utc).strftime('%B %Y')} · Confidential</div>
    </div>""")

    # Summary
    if summary:
        parts.append(f"""<section>
          <div class="label">Executive summary</div>
          <div class="callout">{e(_safe(summary, 4000))}</div>
        </section>""")

    # Metrics
    if metrics_v:
        cells = _to_string_list(metrics_v)
        if not cells:
            cells = [
                f"<strong>{e(d.get('metric', ''))}</strong><br/>{e(d.get('target', d.get('baseline', '')))}"
                for d in _to_dict_list(metrics_v)
            ]
        cells_html = "".join(
            f'<div class="kpi"><h4>Metric {i + 1}</h4><div class="v">{c}</div></div>'
            for i, c in enumerate(cells[:8])
        )
        parts.append(f"""<section>
          <div class="label">{e(_pretty(metrics_k or 'Metrics'))}</div>
          <div class="kpi-grid">{cells_html}</div>
        </section>""")

    # Findings
    if findings_v:
        finds = _to_string_list(findings_v)
        finds_html = "".join(
            f'<div class="find"><div class="num">{i + 1}</div><div>{e(f)}</div></div>'
            for i, f in enumerate(finds[:9])
        )
        parts.append(f"""<section>
          <div class="label">{e(_pretty(findings_k or 'Key findings'))}</div>
          <div class="find-grid">{finds_html}</div>
        </section>""")

    # Recommendations
    if recs_v:
        recs = _to_dict_list(recs_v)
        rec_html_parts: list[str] = []
        for r in recs[:8]:
            t = str(r.get("recommendation") or r.get("title") or r.get("area") or r.get("risk") or "")
            rat = str(r.get("rationale") or r.get("mitigation") or r.get("description") or "")
            pri = str(r.get("priority", "")).lower()
            chip_cls = "chip-high" if "high" in pri else "chip-low" if "low" in pri else "chip-med"
            chip_html = f'<span class="chip {chip_cls}">{e(pri)}</span>' if pri else ""
            actions = _to_string_list(r.get("actions"))
            chips_html = "".join(f"<span>{e(a)}</span>" for a in actions[:6])
            rec_html_parts.append(f"""<div class="rec">
              <div class="top"><h3>{e(t)}</h3>{chip_html}</div>
              {f'<p>{e(rat)}</p>' if rat else ''}
              {f'<div class="chips">{chips_html}</div>' if chips_html else ''}
            </div>""")
        parts.append(f"""<section>
          <div class="label">{e(_pretty(recs_k or 'Recommendations'))}</div>
          {"".join(rec_html_parts)}
        </section>""")

    # Roadmap
    if roadmap_v:
        rows = _to_dict_list(roadmap_v)
        phases_html: list[str] = []
        for i, r in enumerate(rows[:8]):
            phase = str(r.get("phase") or r.get("week_range") or f"Phase {i + 1}")
            tl = str(r.get("timeline") or "")
            acts = _to_string_list(r.get("actions") or r.get("milestones") or r.get("focus_areas"))
            outcome = str(r.get("expected_outcome") or "")
            acts_html = "".join(f"<li>{e(a)}</li>" for a in acts[:6])
            phases_html.append(f"""<div class="phase">
              <h3>{e(phase)}</h3>
              {f'<div class="tl">{e(tl)}</div>' if tl else ''}
              {f'<ul>{acts_html}</ul>' if acts_html else ''}
              {f'<p style="font-size:12px;color:#64748b;margin-top:8px"><em>Outcome:</em> {e(outcome)}</p>' if outcome else ''}
            </div>""")
        parts.append(f"""<section>
          <div class="label">{e(_pretty(roadmap_k or 'Roadmap'))}</div>
          <div class="timeline">{"".join(phases_html)}</div>
        </section>""")

    # Remaining sections
    used = {findings_k, recs_k, metrics_k, roadmap_k} | set(_SUMMARY_KEYS)
    rest = [(k, v) for k, v in source.items() if k not in used and v]
    if rest:
        rest_html_parts: list[str] = []
        for key, value in rest:
            inner: str
            if isinstance(value, str):
                inner = f"<p>{e(_safe(value, 4000))}</p>"
            elif isinstance(value, list):
                items: list[str] = []
                for item in value:
                    if isinstance(item, str):
                        items.append(f"<li>{e(item)}</li>")
                    elif isinstance(item, dict):
                        line = " · ".join(f"<strong>{e(_pretty(k))}:</strong> {e(v)}"
                                          for k, v in list(item.items())[:3] if v)
                        items.append(f"<li>{line}</li>")
                inner = f"<ul>{''.join(items)}</ul>" if items else ""
            elif isinstance(value, dict):
                inner = "".join(
                    f"<p><strong>{e(_pretty(k))}:</strong> {e(_safe(v, 800))}</p>"
                    for k, v in value.items() if v
                )
            else:
                inner = f"<p>{e(_safe(value, 800))}</p>"
            rest_html_parts.append(f'<div class="section-card"><h2>{e(_pretty(key))}</h2>{inner}</div>')
        parts.append(f"""<section>
          <div class="label">Additional sections</div>
          {"".join(rest_html_parts)}
        </section>""")

    parts.append(f"""<footer>Aivora HC Platform · Confidential · {datetime.now(timezone.utc).year}</footer>
</div></body></html>""")

    return "".join(parts).encode("utf-8")
